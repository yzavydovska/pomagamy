"""
Generuje prezentację obrony pracy inżynierskiej — PomagaMY.
Uruchom: python generate_presentation.py
Po dodaniu filmów do folderu videos/ uruchom ponownie — zostaną osadzone w slajdach demo.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
VIDEOS = ROOT / "videos"
OUTPUT = ROOT / "PomagaMY-obrona-pracy-v4.pptx"

# Kolorystyka aplikacji (mobile/src/theme/colors.ts)
PRIMARY = RGBColor(0xB6, 0x7D, 0x2B)
PRIMARY_DARK = RGBColor(0x9A, 0x68, 0x24)
BG = RGBColor(0xF5, 0xF5, 0xF5)
SURFACE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT = RGBColor(0x1A, 0x1A, 0x1A)
TEXT_MUTED = RGBColor(0x6B, 0x65, 0x60)
BORDER = RGBColor(0xE7, 0xE2, 0xDA)
SUCCESS = RGBColor(0x2E, 0x7D, 0x32)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(
    slide,
    left,
    top,
    width,
    height,
    fill: RGBColor,
    line: RGBColor | None = None,
    radius=None,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    if radius and hasattr(shape, "adjustments") and len(shape.adjustments) > 0:
        shape.adjustments[0] = radius
    return shape


def add_textbox(slide, left, top, width, height, text, size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.alignment = align
    return box


def add_bullets(slide, left, top, width, height, items: list[str], size=16, color=TEXT, spacing=1.15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(8)
        p.line_spacing = spacing
    return box


def slide_header(slide, title: str, subtitle: str | None = None, dark=False):
    accent_h = Inches(0.12)
    add_rect(slide, 0, 0, SLIDE_W, accent_h, PRIMARY if not dark else PRIMARY_DARK)
    title_color = TEXT if not dark else WHITE
    sub_color = TEXT_MUTED if not dark else RGBColor(0xF0, 0xE8, 0xDC)
    add_textbox(slide, Inches(0.65), Inches(0.35), Inches(11.5), Inches(0.7), title, size=32, bold=True, color=title_color)
    if subtitle:
        add_textbox(slide, Inches(0.65), Inches(0.95), Inches(11.5), Inches(0.45), subtitle, size=14, color=sub_color)
    add_rect(slide, Inches(0.65), Inches(1.45), Inches(1.2), Inches(0.06), PRIMARY)


def slide_footer(slide, num: int, total: int = 9):
    add_textbox(
        slide,
        Inches(11.8),
        Inches(7.05),
        Inches(1.2),
        Inches(0.3),
        f"{num} / {total}",
        size=11,
        color=TEXT_MUTED,
        align=PP_ALIGN.RIGHT,
    )
    add_textbox(slide, Inches(0.65), Inches(7.05), Inches(4), Inches(0.3), "Pomagamy", size=11, color=TEXT_MUTED)


def build_title_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PRIMARY_DARK)
    add_rect(slide, Inches(9.2), Inches(-0.5), Inches(5), Inches(5), RGBColor(0x8A, 0x5A, 0x18), radius=0.08)
    add_rect(slide, Inches(-1), Inches(5.2), Inches(4.5), Inches(3.5), RGBColor(0xA6, 0x6E, 0x20), radius=0.08)

    add_textbox(slide, Inches(0.9), Inches(1.6), Inches(8), Inches(1.2), "PomagaMY", size=54, bold=True, color=WHITE)
    add_textbox(
        slide,
        Inches(0.92),
        Inches(2.55),
        Inches(9.5),
        Inches(1.1),
        "Aplikacja mobilna do koordynacji współpracy\nmiędzy wolontariuszami a organizacjami non-profit.",
        size=21,
        color=RGBColor(0xFF, 0xF5, 0xE8),
    )
    add_rect(slide, Inches(0.9), Inches(3.85), Inches(2.4), Inches(0.07), WHITE)

    meta = [
        "Praca inżynierska",
        "Autor: [Imię i nazwisko]",
        "Promotor: [Imię i nazwisko promotora]",
        "Kierunek: [np. Informatyka]",
        "Rok akademicki: 2025/2026",
    ]
    y = 4.2
    for line in meta:
        add_textbox(slide, Inches(0.9), Inches(y), Inches(7), Inches(0.35), line, size=15, color=RGBColor(0xF5, 0xE6, 0xD0))
        y += 0.38


def build_problem_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    slide_header(slide, "Problem i cel projektu", "Dlaczego powstała aplikacja PomagaMY?")
    slide_footer(slide, 2)

    cards = [
        (
            "Problem",
            [
                "Dane rozproszone: e-mail, Excel, Messenger, Facebook",
                "Ryzyko pomyłek i przeoczenia zgłoszeń wolontariuszy",
                "Koordynatorzy tracą czas na ręczne pilnowanie wiadomości",
                "Przy małej kadrze NGO każda godzina administracji to mniej realnej pomocy",
            ],
        ),
        (
            "Cel",
            [
                "Jedno miejsce: ogłoszenia, zgłoszenia, profil, powiadomienia",
                "Prosty proces dla wolontariuszy — bez skomplikowanych formalności",
                "Weryfikacja organizacji (NIP, statut) przed publikacją ofert",
                "Mniej biurokracji — więcej czasu na działania społeczne",
            ],
        ),
        (
            "Grupy użytkowników",
            [
                "Wolontariusz — przegląd ofert, zgłoszenia, profil z awatarem",
                "Organizacja — publikacja zadań, rekrutacja, panel zgłoszeń",
                "Administrator — kolejka weryfikacji NGO, skargi, blokada kont",
                "Trzy oddzielne ścieżki UI w jednej aplikacji mobilnej",
            ],
        ),
    ]
    x_positions = [0.65, 4.55, 8.45]
    for i, (head, bullets) in enumerate(cards):
        left = Inches(x_positions[i])
        card_top = Inches(1.75)
        card_h = Inches(4.35)
        add_rect(slide, left, card_top, Inches(3.75), card_h, SURFACE, BORDER, radius=0.06)
        add_rect(slide, left, card_top, Inches(3.75), Inches(0.52), PRIMARY, radius=0.06)
        add_textbox(
            slide,
            left + Inches(0.22),
            card_top + Inches(0.08),
            Inches(3.2),
            Inches(0.38),
            head,
            size=17,
            bold=True,
            color=WHITE,
        )
        add_bullets(
            slide,
            left + Inches(0.22),
            card_top + Inches(0.62),
            Inches(3.3),
            Inches(3.55),
            [f"• {b}" for b in bullets],
            size=13,
            spacing=1.2,
        )

    # Pasek podsumowujący — wypełnia dolną strefę slajdu
    add_rect(slide, Inches(0.65), Inches(6.2), Inches(12.05), Inches(0.72), RGBColor(0xFA, 0xF8, 0xF5), BORDER, radius=0.04)
    add_rect(slide, Inches(0.65), Inches(6.2), Inches(0.08), Inches(0.72), PRIMARY)
    add_textbox(
        slide,
        Inches(0.85),
        Inches(6.32),
        Inches(11.7),
        Inches(0.5),
        "Efekt: scentralizowany przepływ pracy wolontariuszy i organizacji w jednej aplikacji mobilnej PomagaMY (Android + Firebase).",
        size=13,
        bold=True,
        color=PRIMARY_DARK,
    )


def build_architecture_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    slide_header(slide, "Architektura systemu", "React Native (Expo) · Firebase · tryb online i demonstracyjny")
    slide_footer(slide, 3)

    layers = [
        ("Warstwa prezentacji", "Ekrany RN: wolontariusz, organizacja, admin\nReact Navigation · komponenty UI"),
        ("Logika aplikacji", "PomagaMYContext — sesja, ogłoszenia, zgłoszenia, powiadomienia"),
        ("Usługi chmurowe", "Firebase Auth · Cloud Firestore · Storage REST"),
        ("Bezpieczeństwo", "Role użytkowników · Security Rules · weryfikacja organizacji"),
    ]
    y = 1.85
    for i, (title, desc) in enumerate(layers):
        add_rect(slide, Inches(0.65), Inches(y), Inches(0.35), Inches(0.85), PRIMARY)
        add_textbox(slide, Inches(0.65), Inches(y + 0.15), Inches(0.35), Inches(0.5), str(i + 1), size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(slide, Inches(1.15), Inches(y), Inches(11.5), Inches(0.85), SURFACE, BORDER, radius=0.04)
        add_textbox(slide, Inches(1.35), Inches(y + 0.08), Inches(2.5), Inches(0.35), title, size=15, bold=True, color=PRIMARY_DARK)
        add_textbox(slide, Inches(3.9), Inches(y + 0.08), Inches(8.5), Inches(0.7), desc, size=14, color=TEXT)
        y += 1.05

    add_rect(slide, Inches(0.65), Inches(6.15), Inches(12.05), Inches(0.65), RGBColor(0xE8, 0xF5, 0xE9), SUCCESS)
    add_textbox(
        slide,
        Inches(0.85),
        Inches(6.28),
        Inches(11.5),
        Inches(0.4),
        "Widoczność ogłoszeń dla wolontariuszy zależy od orgVerificationStatus === approved (kod + Firestore Rules).",
        size=13,
        color=SUCCESS,
    )


def build_demo_slide(
    prs: Presentation,
    slide_num: int,
    role: str,
    role_color: RGBColor,
    headline: str,
    bullets: list[str],
    video_filename: str,
):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    slide_header(slide, f"Demonstracja — {role}", headline)
    slide_footer(slide, slide_num)

    add_rect(slide, Inches(0.65), Inches(1.75), Inches(4.35), Inches(5.05), SURFACE, BORDER, radius=0.05)
    add_textbox(slide, Inches(0.9), Inches(1.95), Inches(3.8), Inches(0.4), "Kluczowe funkcje", size=17, bold=True, color=PRIMARY_DARK)
    add_bullets(slide, Inches(0.9), Inches(2.45), Inches(3.85), Inches(4.1), bullets, size=14)

    frame_left = Inches(5.25)
    frame_top = Inches(1.75)
    frame_w = Inches(7.45)
    frame_h = Inches(5.05)

    add_rect(slide, frame_left - Inches(0.08), frame_top - Inches(0.08), frame_w + Inches(0.16), frame_h + Inches(0.16), PRIMARY, radius=0.04)
    add_rect(slide, frame_left, frame_top, frame_w, frame_h, RGBColor(0x1A, 0x1A, 0x1A), radius=0.03)

    video_path = VIDEOS / video_filename
    if video_path.exists():
        try:
            slide.shapes.add_movie(
                str(video_path),
                frame_left + Inches(0.15),
                frame_top + Inches(0.15),
                frame_w - Inches(0.3),
                frame_h - Inches(0.3),
                mime_type="video/mp4",
            )
        except Exception:
            _video_placeholder(slide, frame_left, frame_top, frame_w, frame_h, video_filename)
    else:
        _video_placeholder(slide, frame_left, frame_top, frame_w, frame_h, video_filename)

    badge = add_rect(slide, Inches(5.45), Inches(1.95), Inches(2.2), Inches(0.45), role_color, radius=0.2)
    badge.text_frame.paragraphs[0].text = role.upper()
    badge.text_frame.paragraphs[0].font.size = Pt(11)
    badge.text_frame.paragraphs[0].font.bold = True
    badge.text_frame.paragraphs[0].font.color.rgb = WHITE
    badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def _video_placeholder(slide, left, top, width, height, filename: str):
    inner_top = top + Inches(0.15)
    inner_h = height - Inches(0.3)
    add_textbox(
        slide,
        left + Inches(0.4),
        inner_top + Inches(1.5),
        width - Inches(0.8),
        Inches(0.5),
        "▶  MIEJSCE NA NAGRANIE",
        size=22,
        bold=True,
        color=RGBColor(0xB6, 0x7D, 0x2B),
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        left + Inches(0.4),
        inner_top + Inches(2.15),
        width - Inches(0.8),
        Inches(1.2),
        f"Umieść plik:\nvideos/{filename}\n\nNastępnie uruchom:\npython generate_presentation.py",
        size=13,
        color=RGBColor(0xCC, 0xCC, 0xCC),
        align=PP_ALIGN.CENTER,
    )


def build_flow_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    slide_header(slide, "Przepływ biznesowy", "Od ogłoszenia do zakończenia zadania")
    slide_footer(slide, 8)

    steps = [
        "Organizacja publikuje ogłoszenie",
        "Wolontariusz składa zgłoszenie",
        "Organizacja akceptuje kandydata",
        "Realizacja + powiadomienia in-app",
        "Zakończenie → status ogłoszenia",
    ]
    n = len(steps)
    start_x = 0.55
    step_w = 2.35
    for i, step in enumerate(steps):
        x = Inches(start_x + i * step_w)
        add_rect(slide, x, Inches(2.2), Inches(2.05), Inches(1.05), PRIMARY if i % 2 == 0 else SURFACE, PRIMARY if i % 2 else BORDER, radius=0.08)
        tc = WHITE if i % 2 == 0 else PRIMARY_DARK
        add_textbox(slide, x + Inches(0.12), Inches(2.32), Inches(0.4), Inches(0.35), f"{i+1}", size=14, bold=True, color=tc)
        add_textbox(slide, x + Inches(0.12), Inches(2.65), Inches(1.85), Inches(0.55), step, size=12, bold=True, color=tc)
        if i < n - 1:
            add_textbox(slide, x + Inches(2.05), Inches(2.55), Inches(0.35), Inches(0.35), "→", size=20, bold=True, color=PRIMARY)

    add_rect(slide, Inches(0.65), Inches(3.65), Inches(12.05), Inches(2.85), SURFACE, BORDER, radius=0.05)
    add_textbox(slide, Inches(0.9), Inches(3.85), Inches(11.5), Inches(0.35), "Moderacja i zaufanie", size=17, bold=True, color=PRIMARY_DARK)
    add_bullets(
        slide,
        Inches(0.9),
        Inches(4.3),
        Inches(11.3),
        Inches(2.0),
        [
            "Rejestracja organizacji: NIP, KRS, statut → status pending",
            "Administrator weryfikuje dokumenty w dedykowanym panelu",
            "Po zatwierdzeniu ogłoszenia stają się widoczne (visibleToVolunteers: true)",
            "System skarg i blokada kont (accountSuspended) dla naruszeń regulaminu",
        ],
        size=14,
    )


def build_tech_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    slide_header(slide, "Stack technologiczny", "Narzędzia i biblioteki użyte w implementacji")
    slide_footer(slide, 9)

    cols = [
        ("Frontend mobilny", ["React Native 0.81", "Expo SDK 54", "TypeScript", "React Navigation"]),
        ("Backend (BaaS)", ["Firebase Authentication", "Cloud Firestore", "Firebase Storage", "Security Rules"]),
        ("Integracje", ["expo-image-picker", "expo-document-picker", "expo-file-system", "AsyncStorage"]),
    ]
    for i, (head, items) in enumerate(cols):
        left = Inches(0.65 + i * 4.05)
        add_rect(slide, left, Inches(1.85), Inches(3.85), Inches(4.2), SURFACE, BORDER, radius=0.06)
        add_rect(slide, left, Inches(1.85), Inches(3.85), Inches(0.5), PRIMARY, radius=0.06)
        add_textbox(slide, left + Inches(0.2), Inches(1.93), Inches(3.4), Inches(0.35), head, size=16, bold=True, color=WHITE)
        add_bullets(slide, left + Inches(0.25), Inches(2.55), Inches(3.35), Inches(3.2), [f"• {x}" for x in items], size=14)

    add_textbox(
        slide,
        Inches(0.65),
        Inches(6.25),
        Inches(12),
        Inches(0.5),
        "Platforma docelowa: Android · architektura serverless · synchronizacja stanu w czasie rzeczywistym (onSnapshot)",
        size=13,
        color=TEXT_MUTED,
    )


def build_summary_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PRIMARY_DARK)
    slide_footer(slide, 9)

    add_textbox(slide, Inches(0.9), Inches(1.5), Inches(10), Inches(0.8), "Podsumowanie", size=40, bold=True, color=WHITE)
    points = [
        "Zaprojektowano i zaimplementowano aplikację mobilną PomagaMY dla trzech ról użytkowników.",
        "Scentralizowano proces rekrutacji wolontariuszy i weryfikacji organizacji NGO.",
        "Zastosowano Firebase oraz reguły bezpieczeństwa egzekwujące politykę widoczności ofert.",
        "Przygotowano tryb demonstracyjny umożliwiający testy bez połączenia z chmurą.",
    ]
    y = 2.55
    for p in points:
        add_textbox(slide, Inches(0.95), Inches(y), Inches(11), Inches(0.55), f"✓  {p}", size=17, color=RGBColor(0xFF, 0xF5, 0xE8))
        y += 0.72

    add_rect(slide, Inches(0.9), Inches(5.55), Inches(11.5), Inches(0.07), WHITE)
    add_textbox(slide, Inches(0.9), Inches(5.85), Inches(11), Inches(0.8), "Dziękuję za uwagę.\nPytania?", size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_title_slide(prs)
    build_problem_slide(prs)
    build_architecture_slide(prs)

    build_demo_slide(
        prs,
        4,
        "Wolontariusz",
        RGBColor(0x15, 0x65, 0xC0),
        "Przeglądanie ofert, zgłoszenia i profil",
        [
            "• Lista ogłoszeń z filtrowaniem",
            "• Szczegóły oferty i NIP/KRS organizacji",
            "• Zgłoszenie do zadania (submitApplication)",
            "• Powiadomienia i edycja profilu z awatarem",
        ],
        "wolontariusz.mp4",
    )
    build_demo_slide(
        prs,
        5,
        "Organizacja",
        PRIMARY,
        "Publikacja ofert i rekrutacja wolontariuszy",
        [
            "• Rejestracja z NIP i statutem",
            "• Panel ogłoszeń i formularz publikacji",
            "• Akceptacja i zakończenie zgłoszeń",
            "• Status weryfikacji (pending / approved)",
        ],
        "organizacja.mp4",
    )
    build_demo_slide(
        prs,
        6,
        "Administrator",
        RGBColor(0xC6, 0x28, 0x28),
        "Weryfikacja NGO i moderacja systemu",
        [
            "• Kolejka organizacji do zatwierdzenia",
            "• Podgląd statutu i danych rejestrowych",
            "• Obsługa skarg i blokada kont",
            "• Oddzielna nawigacja (RootNavigator)",
        ],
        "admin.mp4",
    )

    build_flow_slide(prs)
    build_tech_slide(prs)
    build_summary_slide(prs)

    prs.save(OUTPUT)
    print(f"Zapisano: {OUTPUT}")
    missing = [f for f in ["wolontariusz.mp4", "organizacja.mp4", "admin.mp4"] if not (VIDEOS / f).exists()]
    if missing:
        print("Brak filmów (placeholdery na slajdach):", ", ".join(missing))
        print(f"Dodaj je do: {VIDEOS}")
    else:
        print("Filmy osadzone w prezentacji.")


if __name__ == "__main__":
    main()
